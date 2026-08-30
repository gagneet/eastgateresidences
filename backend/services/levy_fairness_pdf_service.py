"""
Levy Fairness PDF & PPTX Service

Generates:
  - Professional PDF Audit reports (A4, multi-section)
  - PPTX presentations for AGMs (8+ slides)
  - CSV impact data
"""
import csv
import html as html_lib
import io
from datetime import datetime

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    HRFlowable, KeepTogether
)

from database import db
from services.settings_service import get_general_settings_or_default

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None
    Inches = Pt = Emu = None
    RGBColor = PP_ALIGN = None

# ── Colour palette ──────────────────────────────────────────────────────────────
_DARK = colors.HexColor("#2F4F4F")
_ACCENT = colors.HexColor("#E07A5F")
_GOLD = colors.HexColor("#F2CC8F")
_LIGHT = colors.HexColor("#F8F8F5")
_BORDER = colors.HexColor("#CCCCCC")
_GREEN = colors.HexColor("#2D6A4F")
_RED = colors.HexColor("#C0392B")


def _fmt_aud(n) -> str:
    """Generated function header.

    Function: _fmt_aud
    Path: backend/services/levy_fairness_pdf_service.py

    Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
    """
    if n is None:
        return "—"
    return f"${abs(float(n)):,.0f}"


def _fmt_pct(n, decimals=1) -> str:
    """Generated function header.

    Function: _fmt_pct
    Path: backend/services/levy_fairness_pdf_service.py

    Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
    """
    if n is None:
        return "—"
    v = float(n) * (100 if abs(float(n)) <= 1 else 1)
    sign = "+" if v > 0 else ""
    return f"{sign}{v:.{decimals}f}%"


def _lbfi_band(score) -> str:
    """Generated function header.

    Function: _lbfi_band
    Path: backend/services/levy_fairness_pdf_service.py

    Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
    """
    if score is None:
        return "N/A"
    s = float(score)
    if s >= 98:
        return "Excellent"
    if s >= 94:
        return "Good"
    if s >= 88:
        return "Moderate"
    return "Poor"


# ── PDF Helper: styled table ────────────────────────────────────────────────────

def _make_table(data: list, col_widths=None, header_rows=1) -> Table:
    """Generated function header.

    Function: _make_table
    Path: backend/services/levy_fairness_pdf_service.py

    Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
    """
    t = Table(data, colWidths=col_widths, hAlign="LEFT", repeatRows=header_rows)
    style = [
        ("BACKGROUND", (0, 0), (-1, header_rows - 1), _DARK),
        ("TEXTCOLOR", (0, 0), (-1, header_rows - 1), colors.white),
        ("FONTNAME", (0, 0), (-1, header_rows - 1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.4, _BORDER),
        ("ROWBACKGROUNDS", (0, header_rows), (-1, -1), [colors.white, _LIGHT]),
        ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
        ("ALIGN", (0, 0), (0, -1), "LEFT"),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
    ]
    t.setStyle(TableStyle(style))
    return t


async def generate_agm_report_pdf(building_id: str) -> bytes:
    """Generates a detailed PDF report of the levy fairness analysis."""
    result = await db.levy_fairness_results_v2.find_one({"building_id": building_id})
    if not result:
        return None

    buffer = io.BytesIO()
    PAGE_W, PAGE_H = A4
    doc = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=2.2 * cm, rightMargin=2.2 * cm,
        topMargin=2.2 * cm, bottomMargin=2.2 * cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=18, textColor=_DARK, spaceAfter=4)
    h2_style = ParagraphStyle("h2", fontName="Helvetica-Bold", fontSize=13, textColor=_DARK, spaceBefore=14,
                              spaceAfter=4)
    h3_style = ParagraphStyle("h3", fontName="Helvetica-Bold", fontSize=10, textColor=_DARK, spaceBefore=10,
                              spaceAfter=4)
    body_style = ParagraphStyle("body", fontName="Helvetica", fontSize=9, leading=13, spaceAfter=4)
    small_style = ParagraphStyle("small", fontName="Helvetica", fontSize=7, textColor=colors.grey, spaceAfter=4)
    accent_style = ParagraphStyle("accent", fontName="Helvetica-Bold", fontSize=9, textColor=_ACCENT, spaceAfter=4)

    # ── Gather data ─────────────────────────────────────────────────────────────
    lbfi = result.get("lbfi", {})
    current_score = lbfi.get("current_score", "N/A")
    benefit_score = lbfi.get("benefit_score", 100)
    fairness_gain = lbfi.get("fairness_gain", 0)
    D_value = lbfi.get("D", "N/A")

    lei_score = result.get("lei_score", result.get("lbfi", {}).get("current_score", "N/A"))
    sei_scheme = float(result.get("sei_scheme", 0) or 0)
    total_budget = float(result.get("total_budget", 0) or 0)
    subsidy_map = result.get("subsidy_map", {})
    group_summary = subsidy_map.get("group_summary", [])
    flows = subsidy_map.get("flows", [])
    fac_breakdown = result.get("facility_breakdown", [])
    unit_impact = result.get("unit_impact", [])
    impact_by_group = result.get("impact_by_group", [])
    sim = result.get("simulation", {})
    computed_at = result.get("computed_at", datetime.now().isoformat())

    # ── Elements list ───────────────────────────────────────────────────────────
    E = []

    # ── Cover / Title ───────────────────────────────────────────────────────────
    E.append(Paragraph("East Gate Residences", title_style))
    E.append(Paragraph("Levy Equity Analysis — Committee Report",
                       ParagraphStyle("sub", fontName="Helvetica", fontSize=13, textColor=_ACCENT, spaceAfter=4)))
    E.append(Paragraph(f"Strata Plan {building_id}  |  Generated {datetime.now().strftime('%d %B %Y')}", small_style))
    E.append(HRFlowable(width="100%", thickness=1.5, color=_DARK, spaceAfter=12))

    # ── 1. Executive Summary ────────────────────────────────────────────────────
    E.append(Paragraph("1. Executive Summary", h2_style))
    band = _lbfi_band(current_score)
    net_subsidy = flows[0]["amount"] if flows else 0
    subsidy_dir = f"{flows[0]['from']}s → {flows[0]['to']}s" if flows else "N/A"

    E.append(Paragraph(
        f"This report presents the results of the Levy Benefit Fairness Index (LBFI) analysis for East Gate Residences "
        f"(Strata Plan {building_id}). The analysis compares each lot's proportional levy payment against the facilities "
        f"and services it actually benefits from.",
        body_style
    ))
    E.append(Spacer(1, 6))
    # Summary stats table
    summary_rows = [
        ["Metric", "Value", "Benchmark"],
        ["LBFI Score (current UE model)", f"{current_score} / 100", "100 = perfect alignment"],
        ["LBFI Rating", band, "Good ≥ 94, Excellent ≥ 98"],
        ["Gini-Gap (D)", str(D_value), "0 = no cross-subsidy"],
        ["Benefit-Based Target Score", f"{benefit_score} / 100", "Theoretical maximum"],
        ["Potential Fairness Gain", f"+{fairness_gain} pts", "If benefit model adopted"],
        ["Annual Cross-Subsidy Flow", _fmt_aud(net_subsidy), f"{subsidy_dir}"],
        ["Annual Subsidy Distortion", f"{sei_scheme * 100:.1f}%", "< 5% = acceptable"],
        ["Total Scheme Levy (FY2026)", _fmt_aud(total_budget), "All 87 lots combined"],
        ["Special Levy Probability (10yr)", f"{sim.get('special_levy_probability', 0):.1f}%", "< 10% = low risk"],
    ]
    col_w = [PAGE_W * 0.38, PAGE_W * 0.25, PAGE_W * 0.28]
    E.append(_make_table(summary_rows, col_widths=col_w))
    E.append(Spacer(1, 8))

    E.append(Paragraph(
        f"Under the current Unit of Entitlement (UOE) model all 87 owners pay in proportion to their UOE. "
        f"East Gate has <b>70 Apartments</b> (Tower) and <b>17 Townhouses</b> that use fundamentally different facilities. "
        f"Apartments use lifts, internal corridor cleaning, tower intercom and basement systems. Townhouses have separate "
        f"perimeter fencing and drainage. This creates an annual cross-subsidy of approximately <b>{_fmt_aud(net_subsidy)}</b> "
        f"flowing from {subsidy_dir} each year.",
        body_style
    ))

    # ── 2. LBFI Gauge Summary ───────────────────────────────────────────────────
    E.append(Paragraph("2. Fairness Score Breakdown (LBFI)", h2_style))
    E.append(Paragraph(
        "LBFI = 100 × (1 − D) where D = ½ × Σ|p_i − b_i| is the Gini-style gap between each lot's payment share (p_i) "
        "and benefit share (b_i). A score of 100 means every lot pays exactly proportional to the facilities it uses.",
        body_style
    ))
    gauge_rows = [
        ["Model", "LBFI Score", "Rating", "Annual Total Levy", "Avg / Lot"],
        ["Current (UE-only)", f"{current_score}", band, _fmt_aud(total_budget),
         _fmt_aud(total_budget / 87 if total_budget else 0)],
        ["Benefit-Based Target", f"{benefit_score}", "Excellent", _fmt_aud(total_budget),
         _fmt_aud(total_budget / 87 if total_budget else 0)],
    ]
    E.append(
        _make_table(gauge_rows, col_widths=[PAGE_W * 0.32, PAGE_W * 0.12, PAGE_W * 0.14, PAGE_W * 0.2, PAGE_W * 0.13]))
    E.append(Paragraph(
        "Note: the total levy collected is <b>identical</b> in both models. The benefit-based model only redistributes "
        "cost allocation — it does not change the overall scheme budget.",
        small_style
    ))

    # ── 3. Subsidy Flow Analysis ────────────────────────────────────────────────
    E.append(Paragraph("3. Cross-Subsidy Analysis", h2_style))
    if group_summary:
        E.append(Paragraph(
            "The table below shows each ownership group's current levy total versus what they would pay under a pure "
            "benefit-based model, and the resulting net subsidy position.",
            body_style
        ))
        grp_rows = [["Group", "Lots", "UOE", "Current Total Levy", "Benefit-Based Levy", "Annual Change", "Role"]]
        for g in group_summary:
            delta = (g.get("benefit_total", 0) or 0) - (g.get("current_total", 0) or 0)
            grp_rows.append([
                f"{g.get('group', '?')}s",
                str(g.get("count", "—")),
                str(g.get("ue", "—")),
                _fmt_aud(g.get("current_total")),
                _fmt_aud(g.get("benefit_total")),
                ("+" if delta >= 0 else "") + _fmt_aud(abs(delta)),
                g.get("role", "Neutral"),
            ])
        E.append(_make_table(grp_rows,
                             col_widths=[PAGE_W * 0.15, PAGE_W * 0.07, PAGE_W * 0.08, PAGE_W * 0.16, PAGE_W * 0.16,
                                         PAGE_W * 0.13, PAGE_W * 0.13]))
    if flows:
        E.append(Spacer(1, 6))
        E.append(Paragraph("Annual Subsidy Flows:", h3_style))
        for fl in flows:
            E.append(Paragraph(
                f"• <b>{fl.get('from', '')}s → {fl.get('to', '')}s:</b> {_fmt_aud(fl.get('amount'))} per year "
                f"(Townhouses currently contribute toward apartment-only facilities such as lifts, intercom, and corridor cleaning).",
                body_style
            ))

    # ── 4. Facility Cost Allocation ─────────────────────────────────────────────
    E.append(Paragraph("4. Facility Cost Allocation Centres", h2_style))
    E.append(Paragraph(
        "The three-tier model assigns each facility to the groups that directly benefit from it: "
        "Global (all lots), Apartment-Only, or Townhouse-Only.",
        body_style
    ))
    if fac_breakdown:
        fac_rows = [["Facility", "Annual Cost", "Benefit Group", "Cost Tier"]]
        for f in fac_breakdown:
            fac_rows.append([
                f.get("facility_name", f.get("name", "—")),
                _fmt_aud(f.get("annual_cost")),
                f.get("benefit_group_id", f.get("group", "—")),
                f.get("tier", "global").capitalize(),
            ])
        E.append(_make_table(fac_rows, col_widths=[PAGE_W * 0.38, PAGE_W * 0.15, PAGE_W * 0.22, PAGE_W * 0.15]))
    else:
        E.append(Paragraph("Facility breakdown data not available. Please regenerate the analysis.", body_style))

    # ── 5. Predictive Risk ───────────────────────────────────────────────────────
    if sim:
        E.append(Paragraph("5. Predictive Financial Risk (Monte Carlo Simulation)", h2_style))
        E.append(Paragraph(
            "Based on 1,000 Monte Carlo simulations modelling contribution rates, capital expenditure timing, "
            "and inflation variance over a 10-year horizon:",
            body_style
        ))
        risk_rows = [
            ["Risk Metric", "Result", "Interpretation"],
            ["Special Levy Probability (10yr)", f"{sim.get('special_levy_probability', 0):.1f}%",
             "Probability that a one-off special levy is required within 10 years"],
            ["P50 — Median Reserve Balance", _fmt_aud(sim.get("p50")),
             "50th percentile projected reserve at end of horizon"],
            ["P75 — Upper Reserve Balance", _fmt_aud(sim.get("p75")),
             "75th percentile — reserve if levy income performs above median"],
            ["P90 — Risk Level Reserve", _fmt_aud(sim.get("p90")),
             "Reserve at 90th percentile scenario"],
            ["P95 — Crisis Level Reserve", _fmt_aud(sim.get("p95")),
             "Reserve at 95th percentile — stress scenario"],
        ]
        E.append(_make_table(risk_rows, col_widths=[PAGE_W * 0.38, PAGE_W * 0.15, PAGE_W * 0.38]))

    # ── 6. Per-Unit Impact Table (first 50 rows) ─────────────────────────────────
    if unit_impact:
        E.append(Paragraph("6. Per-Lot Impact Summary", h2_style))
        E.append(Paragraph(
            f"The table below shows the levy impact per lot under the benefit-based model. "
            f"Lots with a negative change (in green) would pay less; those with a positive change would pay more. "
            f"Total change across all lots is $0 — only the distribution changes. "
            f"Showing first {min(50, len(unit_impact))} of {len(unit_impact)} lots.",
            body_style
        ))
        unit_rows = [["Unit", "Type", "Owner", "Current Levy", "Fair Levy", "Proposed Levy", "Change", "SEI"]]
        for u in unit_impact[:50]:
            change = float(u.get("change") or u.get("proposed_levy", 0) - u.get("current_levy", 0))
            unit_rows.append([
                u.get("unit_number", "—"),
                u.get("unit_type", "—").capitalize(),
                (u.get("owner_name") or "—")[:20],
                _fmt_aud(u.get("current_levy")),
                _fmt_aud(u.get("fair_levy")),
                _fmt_aud(u.get("proposed_levy")),
                ("+" if change >= 0 else "") + _fmt_aud(abs(change)),
                f"{float(u.get('sei', 0) or 0) * 100:.1f}%",
            ])
        t = _make_table(unit_rows, col_widths=[PAGE_W * 0.09, PAGE_W * 0.09, PAGE_W * 0.18, PAGE_W * 0.1, PAGE_W * 0.1,
                                               PAGE_W * 0.12, PAGE_W * 0.1, PAGE_W * 0.09])
        # Colour rows with negative change green, positive red
        style_extra = []
        for i, u in enumerate(unit_impact[:50], start=1):
            change = float(u.get("change") or 0)
            if change < -0.5:
                style_extra.append(("TEXTCOLOR", (6, i), (6, i), _GREEN))
            elif change > 0.5:
                style_extra.append(("TEXTCOLOR", (6, i), (6, i), _RED))
        t._argW = t._argW  # force re-validate
        E.append(t)
        if len(unit_impact) > 50:
            E.append(Paragraph(
                f"Full per-lot data for all {len(unit_impact)} lots available in the CSV export (Raw Data download).",
                small_style
            ))

    # ── 7. Legal Pathways ───────────────────────────────────────────────────────
    E.append(Paragraph("7. Legal Pathway Options", h2_style))
    E.append(Paragraph(
        "Three legal mechanisms exist under ACT and NSW strata legislation to implement benefit-based levy allocation:",
        body_style
    ))
    legal_rows = [
        ["Option", "Mechanism", "Complexity", "Vote Required", "Timeline"],
        ["1 — Recommended",
         "Exclusive Use By-Law\n(ACT UTM Act 2011 s.78 / NSW SSMA 2015 s.107)\nAssigns apartment-specific facility costs to Apartment lots via registered by-law.",
         "Low–Medium", "Special Resolution\n(75% of vote value)", "3–6 months"],
        ["2 — Structural",
         "Stratum Subdivision + BMC\nSplits the scheme into two strata plans (Apartments & Townhouses) with a Building Management Committee.",
         "Very High", "Tribunal/Court approval\n+ subdivision process", "12–24 months"],
        ["3 — Last Resort",
         "UOE Re-apportionment\nFormally reapportions Unit of Entitlement values via ACT Civil & Administrative Tribunal.",
         "Extreme", "ACAT approval\n+ unanimous consent", "18–36 months"],
    ]
    legal_t = Table(legal_rows, colWidths=[PAGE_W * 0.12, PAGE_W * 0.38, PAGE_W * 0.13, PAGE_W * 0.16, PAGE_W * 0.12],
                    hAlign="LEFT")
    legal_t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), _DARK),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 7),
        ("GRID", (0, 0), (-1, -1), 0.4, _BORDER),
        ("BACKGROUND", (0, 1), (-1, 1), colors.HexColor("#EAF5EA")),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
    ]))
    E.append(legal_t)

    # ── 8. Recommended Action Plan ──────────────────────────────────────────────
    E.append(Paragraph("8. Recommended Action Plan for AGM", h2_style))
    steps = [
        ("Immediate — Before AGM",
         "Engage a strata lawyer to draft the Exclusive Use By-Law. "
         "Frame the by-law to assign cost responsibility for lifts, tower intercom, corridor cleaning, "
         "and basement systems exclusively to Apartment lot owners."),
        ("AGM Resolution",
         f"Present this report and LBFI analysis to all owners. "
         f"Propose a Special Resolution requiring 75% of total vote value. "
         f"Current LBFI score of {current_score}/100 provides clear evidence of cross-subsidy."),
        ("Year 1 Transition",
         f"Apply a maximum {_fmt_pct(0.10)} annual levy increase cap per lot. "
         f"Begin transitioning levies toward the benefit-based model progressively."),
        ("Years 2–3 Transition",
         "Continue phased transition. Most lots will reach the benefit-based target within 3 years "
         "under a 10% annual cap. Monitor and adjust."),
        ("Ongoing",
         "Update the Digital Twin (facility register) annually. "
         "Rerun the LBFI engine at each budget cycle to confirm alignment."),
    ]
    step_rows = [["Step", "Action"]]
    for s in steps:
        step_rows.append([s[0], s[1]])
    E.append(_make_table(step_rows, col_widths=[PAGE_W * 0.22, PAGE_W * 0.69]))

    # ── Disclaimer ──────────────────────────────────────────────────────────────
    E.append(Spacer(1, 14))
    E.append(HRFlowable(width="100%", thickness=0.5, color=_BORDER, spaceAfter=6))
    E.append(Paragraph(
        f"<i>This report is generated by the East Gate Residences Strata Management Platform. "
        f"Analysis computed {computed_at[:10] if computed_at else 'unknown'}. "
        f"Levy figures are based on FY2026 actual levy data from the strata manager's records. "
        f"Legal pathway information is indicative only — consult a licensed strata lawyer before proceeding.</i>",
        small_style
    ))

    doc.build(E)
    return buffer.getvalue()


async def generate_agm_presentation(building_id: str) -> bytes:
    """Generates a comprehensive PPTX presentation for the AGM (8 slides)."""
    if not PPTX_AVAILABLE:
        return None

    result = await db.levy_fairness_results_v2.find_one({"building_id": building_id})
    if not result:
        return None

    lbfi = result.get("lbfi", {})
    current_score = lbfi.get("current_score", "N/A")
    benefit_score = lbfi.get("benefit_score", 100)
    fairness_gain = lbfi.get("fairness_gain", 0)
    lei_score = result.get("lei_score", current_score)
    sei_scheme = float(result.get("sei_scheme", 0) or 0)
    total_budget = float(result.get("total_budget", 0) or 0)
    subsidy_map = result.get("subsidy_map", {})
    group_summary = subsidy_map.get("group_summary", [])
    flows = subsidy_map.get("flows", [])
    fac_breakdown = result.get("facility_breakdown", [])
    unit_impact = result.get("unit_impact", [])
    impact_by_group = result.get("impact_by_group", [])
    sim = result.get("simulation", {})
    computed_date = datetime.now().strftime("%d %B %Y")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_HEX = RGBColor(0x2F, 0x4F, 0x4F)
    ACCENT_HEX = RGBColor(0xE0, 0x7A, 0x5F)
    GOLD_HEX = RGBColor(0xF2, 0xCC, 0x8F)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GREY = RGBColor(0xF4, 0xF4, 0xF0)

    def _blank_slide():
        """Generated function header.

        Function: _blank_slide
        Path: backend/services/levy_fairness_pdf_service.py

        Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
        """
        return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    def _add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None):
        """Generated function header.

        Function: _add_rect
        Path: backend/services/levy_fairness_pdf_service.py

        Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
        """
        from pptx.util import Emu
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            shape.fill.background()
        if line_rgb:
            shape.line.color.rgb = line_rgb
        else:
            shape.line.fill.background()
        return shape

    def _add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                      color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
        """Generated function header.

        Function: _add_text_box
        Path: backend/services/levy_fairness_pdf_service.py

        Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
        """
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return tb

    def _add_bullet_box(slide, left, top, width, height, title, bullets, title_size=14, bullet_size=11):
        """Generated function header.

        Function: _add_bullet_box
        Path: backend/services/levy_fairness_pdf_service.py

        Note: Generated inventory header. Replace or expand this with reviewed business-purpose documentation before relying on it as source commentary.
        """
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        # Title paragraph
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = DARK_HEX
        # Bullet paragraphs
        for b in bullets:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(2)
            r2 = p2.add_run()
            r2.text = f"• {b}"
            r2.font.size = Pt(bullet_size)

    net_subsidy = flows[0]["amount"] if flows else 0
    subsidy_from = flows[0]["from"] if flows else "Townhouses"
    subsidy_to = flows[0]["to"] if flows else "Apartments"
    band = _lbfi_band(current_score)

    # Fetch building settings for branding
    b_settings = await get_general_settings_or_default(building_id, {"_id": 0})
    building_name = b_settings.get("building_name", "East Gate Residences")
    plan_number = b_settings.get("plan_number", "13195")

    # ── Slide 1: Title ──────────────────────────────────────────────────────────
    s1 = _blank_slide()
    _add_rect(s1, 0, 0, 13.33, 7.5, fill_rgb=DARK_HEX)
    _add_rect(s1, 0, 5.2, 13.33, 2.3, fill_rgb=ACCENT_HEX)
    _add_text_box(s1, 0.5, 0.4, 12, 1.0, html_lib.escape(building_name), font_size=28, bold=True, color=WHITE,
                  align=PP_ALIGN.LEFT)
    _add_text_box(s1, 0.5, 1.3, 12, 0.8, "Levy Fairness Equity Analysis", font_size=22, bold=False, color=GOLD_HEX,
                  align=PP_ALIGN.LEFT)
    _add_text_box(s1, 0.5, 2.0, 12, 0.6, f"Strata Plan {html_lib.escape(plan_number)}  |  {computed_date}",
                  font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    _add_text_box(s1, 0.5, 5.4, 12, 0.7, "Annual General Meeting — Levy Benefit Fairness Index (LBFI) Report",
                  font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_text_box(s1, 0.5, 6.1, 12, 0.6, "Prepared by the Executive Committee — Confidential, For Owner Distribution",
                  font_size=10, italic=True, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    # ── Slide 2: Agenda ─────────────────────────────────────────────────────────
    s2 = _blank_slide()
    _add_rect(s2, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
    _add_text_box(s2, 0.4, 0.15, 12, 0.8, "Agenda", font_size=22, bold=True, color=WHITE)
    agenda_items = [
        "1. What is the LBFI and why does it matter?",
        "2. Current fairness score and cross-subsidy flows",
        "3. Facility cost breakdown by benefit group",
        "4. Per-lot impact of moving to a benefit-based model",
        "5. Predictive financial risk (Monte Carlo)",
        "6. Legal pathway options",
        "7. Recommended action plan and resolution",
    ]
    _add_bullet_box(s2, 1.0, 1.3, 11, 5.5, "", agenda_items, bullet_size=14)

    # ── Slide 3: What is LBFI? ─────────────────────────────────────────────────
    s3 = _blank_slide()
    _add_rect(s3, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
    _add_text_box(s3, 0.4, 0.15, 12, 0.8, "What is the LBFI?", font_size=22, bold=True, color=WHITE)
    _add_text_box(s3, 0.5, 1.2, 7.5, 5.5,
                  "The Levy Benefit Fairness Index measures how well each lot's levy payment aligns with the facilities it actually uses.\n\n"
                  "LBFI = 100 × (1 − D)\n"
                  "D = Gini-style gap between payment share and benefit share\n\n"
                  "Score 100 = every lot pays exactly for what it uses\n"
                  "Score < 88 = significant cross-subsidies present\n\n"
                  "East Gate has 70 Apartments and 17 Townhouses. Apartments use lifts, tower intercom, corridor cleaning, and "
                  "basement parking systems. Townhouses use separate perimeter fencing and drainage. Under the current flat UOE "
                  "model, Townhouses partially fund facilities they never access.",
                  font_size=12, color=DARK_HEX)
    # Score box
    _add_rect(s3, 8.8, 1.5, 4.0, 2.5, fill_rgb=LIGHT_GREY)
    _add_text_box(s3, 9.0, 1.6, 3.6, 0.5, "Current LBFI Score", font_size=12, bold=True, color=DARK_HEX,
                  align=PP_ALIGN.CENTER)
    _add_text_box(s3, 9.0, 2.0, 3.6, 1.2, str(current_score), font_size=48, bold=True,
                  color=DARK_HEX if float(current_score or 0) >= 88 else ACCENT_HEX, align=PP_ALIGN.CENTER)
    _add_text_box(s3, 9.0, 3.0, 3.6, 0.5, f"Rating: {band}", font_size=13, bold=True,
                  color=RGBColor(0x2D, 0x6A, 0x4F) if band in ("Excellent", "Good") else ACCENT_HEX,
                  align=PP_ALIGN.CENTER)
    _add_text_box(s3, 9.0, 4.2, 3.6, 0.8,
                  f"Benefit-based target: {benefit_score}/100\n(+{fairness_gain} pts gain achievable)",
                  font_size=11, color=DARK_HEX, align=PP_ALIGN.CENTER)

    # ── Slide 4: Cross-Subsidy Analysis ────────────────────────────────────────
    s4 = _blank_slide()
    _add_rect(s4, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
    _add_text_box(s4, 0.4, 0.15, 12, 0.8, "Cross-Subsidy Analysis", font_size=22, bold=True, color=WHITE)
    # Group cards
    x_positions = [0.5, 7.0]
    for idx, g in enumerate(group_summary[:2]):
        x = x_positions[idx]
        w = 5.8
        delta = (g.get("benefit_total", 0) or 0) - (g.get("current_total", 0) or 0)
        role = g.get("role", "Neutral")
        bg = RGBColor(0xEA, 0xF5, 0xEA) if role == "Recipient" else RGBColor(0xFD, 0xEC, 0xEC)
        _add_rect(s4, x, 1.3, w, 5.5, fill_rgb=bg)
        _add_text_box(s4, x + 0.2, 1.4, w - 0.4, 0.6,
                      f"{g.get('group', '?')}s  ({g.get('count', 0)} lots, UOE {g.get('ue', 0)})",
                      font_size=14, bold=True, color=DARK_HEX)
        _add_text_box(s4, x + 0.2, 2.0, w - 0.4, 0.5, f"Current Levy: {_fmt_aud(g.get('current_total'))}", font_size=12,
                      color=DARK_HEX)
        _add_text_box(s4, x + 0.2, 2.5, w - 0.4, 0.5, f"Benefit-Based Levy: {_fmt_aud(g.get('benefit_total'))}",
                      font_size=12, color=DARK_HEX)
        sign = "+" if delta >= 0 else ""
        delta_color = ACCENT_HEX if delta > 0 else RGBColor(0x2D, 0x6A, 0x4F)
        _add_text_box(s4, x + 0.2, 3.0, w - 0.4, 0.7,
                      f"Annual change: {sign}{_fmt_aud(abs(delta))}", font_size=13, bold=True, color=delta_color)
        _add_text_box(s4, x + 0.2, 3.7, w - 0.4, 0.5, f"Role: {role}", font_size=12, bold=True, color=DARK_HEX)
        avg_current = (g.get("current_total", 0) or 0) / (g.get("count", 1) or 1)
        avg_benefit = (g.get("benefit_total", 0) or 0) / (g.get("count", 1) or 1)
        _add_text_box(s4, x + 0.2, 4.3, w - 0.4, 0.5, f"Avg per lot (current): {_fmt_aud(avg_current)}", font_size=11,
                      color=DARK_HEX)
        _add_text_box(s4, x + 0.2, 4.8, w - 0.4, 0.5, f"Avg per lot (target): {_fmt_aud(avg_benefit)}", font_size=11,
                      color=DARK_HEX)
    if flows:
        _add_text_box(s4, 0.4, 6.7, 12, 0.5,
                      f"Annual subsidy flow: {subsidy_from}s currently contribute {_fmt_aud(net_subsidy)}/yr toward apartment-only facilities.",
                      font_size=11, italic=True, color=DARK_HEX)

    # ── Slide 5: Facility Breakdown ─────────────────────────────────────────────
    s5 = _blank_slide()
    _add_rect(s5, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
    _add_text_box(s5, 0.4, 0.15, 12, 0.8, "Facility Cost Allocation by Benefit Tier", font_size=22, bold=True,
                  color=WHITE)
    _add_text_box(s5, 0.4, 1.2, 12.5, 0.5,
                  "Three-tier model: Global (all lots) | Apartment-Only | Townhouse-Only. Each facility assigned to the group that benefits.",
                  font_size=11, color=DARK_HEX)
    if fac_breakdown:
        bullets = []
        for f in fac_breakdown[:12]:
            tier = f.get("tier", f.get("benefit_group_id", "global")).capitalize()
            bullets.append(
                f"{f.get('facility_name', f.get('name', '—'))}  —  {_fmt_aud(f.get('annual_cost'))}  [{tier}]")
        _add_bullet_box(s5, 0.5, 1.8, 12.5, 5.2, "Facility  |  Annual Cost  |  Benefit Tier", bullets, bullet_size=11)
        global_total = sum(float(f.get("annual_cost", 0) or 0) for f in fac_breakdown if
                           "global" in str(f.get("tier", "global")).lower())
        apt_total = sum(
            float(f.get("annual_cost", 0) or 0) for f in fac_breakdown if "apart" in str(f.get("tier", "")).lower())
        th_total = sum(
            float(f.get("annual_cost", 0) or 0) for f in fac_breakdown if "town" in str(f.get("tier", "")).lower())
        _add_text_box(s5, 0.5, 6.8, 12.5, 0.5,
                      f"Totals — Global: {_fmt_aud(global_total)}  |  Apartment-Only: {_fmt_aud(apt_total)}  |  Townhouse-Only: {_fmt_aud(th_total)}",
                      font_size=10, bold=True, color=DARK_HEX)
    else:
        _add_text_box(s5, 0.5, 1.8, 12, 1,
                      "Facility data not available. Regenerate analysis from the Intelligence Hub.", font_size=12,
                      color=DARK_HEX)

    # ── Slide 6: Per-Lot Impact Highlights ─────────────────────────────────────
    s6 = _blank_slide()
    _add_rect(s6, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
    _add_text_box(s6, 0.4, 0.15, 12, 0.8, "Per-Lot Impact: Who Pays More / Less?", font_size=22, bold=True, color=WHITE)
    lots_lower = sum(1 for u in unit_impact if (u.get("change") or 0) < -0.5)
    lots_higher = sum(1 for u in unit_impact if (u.get("change") or 0) > 0.5)
    lots_same = len(unit_impact) - lots_lower - lots_higher
    avg_saving = (sum(u.get("change", 0) or 0 for u in unit_impact if
                      (u.get("change") or 0) < -0.5) / lots_lower) if lots_lower else 0
    avg_increase = (sum(u.get("change", 0) or 0 for u in unit_impact if
                        (u.get("change") or 0) > 0.5) / lots_higher) if lots_higher else 0
    # Stats boxes
    stats = [
        (0.5, "Lots with Lower Levy", str(lots_lower), f"Avg saving: {_fmt_aud(abs(avg_saving))}/yr",
         RGBColor(0xEA, 0xF5, 0xEA)),
        (4.8, "Lots Unchanged", str(lots_same), "Within $0.50 of current", LIGHT_GREY),
        (9.1, "Lots with Higher Levy", str(lots_higher), f"Avg increase: {_fmt_aud(avg_increase)}/yr",
         RGBColor(0xFD, 0xEC, 0xEC)),
    ]
    for x, lbl, val, sub, bg in stats:
        _add_rect(s6, x, 1.3, 3.7, 2.5, fill_rgb=bg)
        _add_text_box(s6, x + 0.2, 1.4, 3.3, 0.6, lbl, font_size=11, bold=True, color=DARK_HEX, align=PP_ALIGN.CENTER)
        _add_text_box(s6, x + 0.2, 1.9, 3.3, 0.9, val, font_size=32, bold=True, color=DARK_HEX, align=PP_ALIGN.CENTER)
        _add_text_box(s6, x + 0.2, 2.8, 3.3, 0.6, sub, font_size=10, color=DARK_HEX, align=PP_ALIGN.CENTER)
    _add_text_box(s6, 0.5, 4.1, 12.3, 0.6,
                  "Important: Total scheme levy is UNCHANGED. The benefit-based model only redistributes costs — it does not increase the overall budget.",
                  font_size=12, bold=True, color=ACCENT_HEX, align=PP_ALIGN.CENTER)
    # Top savers and top movers
    sorted_impact = sorted(unit_impact, key=lambda u: u.get("change", 0) or 0)
    savers = sorted_impact[:5]
    movers = sorted_impact[-5:][::-1]
    _add_bullet_box(s6, 0.5, 4.9, 6.0, 2.3, "Largest Savings (per lot/yr):",
                    [
                        f"{u.get('unit_number', '?')} ({u.get('unit_type', '?')}): {_fmt_aud(abs(u.get('change', 0)))} saving"
                        for u in savers],
                    title_size=11, bullet_size=10)
    _add_bullet_box(s6, 7.0, 4.9, 6.0, 2.3, "Largest Increases (per lot/yr):",
                    [
                        f"{u.get('unit_number', '?')} ({u.get('unit_type', '?')}): {_fmt_aud(abs(u.get('change', 0)))} increase"
                        for u in movers],
                    title_size=11, bullet_size=10)

    # ── Slide 7: Financial Risk ─────────────────────────────────────────────────
    if sim:
        s7 = _blank_slide()
        _add_rect(s7, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
        _add_text_box(s7, 0.4, 0.15, 12, 0.8, "Predictive Financial Risk (Monte Carlo)", font_size=22, bold=True,
                      color=WHITE)
        prob = float(sim.get("special_levy_probability", 0) or 0)
        prob_color = ACCENT_HEX if prob > 20 else RGBColor(0x2D, 0x6A, 0x4F)
        _add_rect(s7, 0.5, 1.3, 4.5, 5.5, fill_rgb=LIGHT_GREY)
        _add_text_box(s7, 0.7, 1.5, 4.1, 0.6, "Special Levy Probability (10yr)", font_size=12, bold=True,
                      color=DARK_HEX, align=PP_ALIGN.CENTER)
        _add_text_box(s7, 0.7, 2.0, 4.1, 1.4, f"{prob:.1f}%", font_size=52, bold=True, color=prob_color,
                      align=PP_ALIGN.CENTER)
        _add_text_box(s7, 0.7, 3.3, 4.1, 0.5,
                      "Low Risk (< 10%)" if prob < 10 else "Elevated Risk (10–30%)" if prob < 30 else "High Risk (> 30%)",
                      font_size=12, bold=True, color=prob_color, align=PP_ALIGN.CENTER)
        _add_text_box(s7, 0.7, 4.0, 4.1, 2.0,
                      f"Based on 1,000 simulations modelling contribution rates, capital expenditure timing, and inflation over 10 years.",
                      font_size=10, color=DARK_HEX, align=PP_ALIGN.CENTER)
        reserve_bullets = [
            f"P50 (Median Reserve): {_fmt_aud(sim.get('p50'))}",
            f"P75 (Above-median): {_fmt_aud(sim.get('p75'))}",
            f"P90 (Risk level): {_fmt_aud(sim.get('p90'))}",
            f"P95 (Crisis level): {_fmt_aud(sim.get('p95'))}",
        ]
        _add_bullet_box(s7, 5.5, 1.3, 7.5, 5.5, "Projected Reserve Balance Scenarios:", reserve_bullets, bullet_size=13)
    else:
        s7 = _blank_slide()
        _add_rect(s7, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
        _add_text_box(s7, 0.4, 0.15, 12, 0.8, "Predictive Financial Risk", font_size=22, bold=True, color=WHITE)
        _add_text_box(s7, 0.5, 1.5, 12, 1,
                      "Monte Carlo simulation data not available. Run 'Regenerate Model' to compute.", font_size=14,
                      color=DARK_HEX)

    # ── Slide 8: Legal Pathways + Action Plan ───────────────────────────────────
    s8 = _blank_slide()
    _add_rect(s8, 0, 0, 13.33, 1.1, fill_rgb=DARK_HEX)
    _add_text_box(s8, 0.4, 0.15, 12, 0.8, "Legal Pathways & Recommended Action Plan", font_size=22, bold=True,
                  color=WHITE)
    legal_bullets = [
        "OPTION 1 — RECOMMENDED: Exclusive Use By-Law (ACT UTM Act 2011 s.78)",
        "  Assigns apartment-specific facility costs to Apartment lots. Low–Medium complexity.",
        "  Requires Special Resolution (75% of total vote value). Timeline: 3–6 months.",
        "",
        "OPTION 2: Stratum Subdivision + BMC — Very High complexity, 12–24 months.",
        "OPTION 3: UOE Re-apportionment — Last resort, requires ACAT approval, 18–36 months.",
    ]
    _add_bullet_box(s8, 0.5, 1.2, 12.5, 3.2, "Legal Mechanisms:", legal_bullets, title_size=13, bullet_size=11)
    action_bullets = [
        "1. Engage strata lawyer to draft Exclusive Use By-Law BEFORE AGM.",
        "2. Present LBFI report at AGM — propose Special Resolution.",
        "3. Apply 10% annual transition cap in Year 1 levy notices.",
        "4. Full transition to benefit-based model within 3 years.",
        "5. Rerun LBFI annually at each budget cycle.",
    ]
    _add_bullet_box(s8, 0.5, 4.5, 12.5, 2.8, "Action Plan:", action_bullets, title_size=13, bullet_size=11)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


async def generate_impact_csv(building_id: str) -> str:
    """Generates CSV content for per-unit impacts."""
    result = await db.levy_fairness_results_v2.find_one({"building_id": building_id})
    if not result or "unit_impact" not in result:
        return ""

    output = io.StringIO()
    writer = csv.DictWriter(
        output,
        fieldnames=["unit_number", "unit_type", "owner_name", "current_levy", "fair_levy", "proposed_levy", "change",
                    "sei"],
        extrasaction="ignore"
    )
    writer.writeheader()
    writer.writerows(result["unit_impact"])
    return output.getvalue()
